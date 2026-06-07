# services/file_service.py

import os
import re
import logging
from typing import Optional

logger = logging.getLogger(__name__)

# Try importing optional libraries
try:
    import pdfplumber
    HAS_PDFPLUMBER = True
except ImportError:
    HAS_PDFPLUMBER = False

try:
    import pytesseract
    from PIL import Image
    import pdf2image
    HAS_OCR = True
except ImportError:
    HAS_OCR = False

try:
    from docx import Document
    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False

try:
    import openpyxl
    HAS_OPENPYXL = True
except ImportError:
    HAS_OPENPYXL = False


class FileService:
    """Service for extracting text from various file formats."""
    
    ALLOWED_EXTENSIONS = {'.pdf', '.docx', '.txt', '.pptx', '.xlsx', '.py', '.js', '.html', '.css', '.json', '.md'}
    MAX_FILE_SIZE = 50 * 1024 * 1024  # 50MB
    
    # کلمات فارسی رایبر برای تشخیص جهت صحیح
    COMMON_PERSIAN_WORDS = {
        'است', 'هست', 'می', 'شود', 'کرد', 'بود', 'این', 'آن', 'با', 'از', 'به', 'در',
        'که', 'را', 'و', 'یا', 'برای', 'تا', 'نیز', 'هم', 'ما', 'شما', 'او', 'آنها',
        'توسعه', 'برنامه', 'سیستم', 'داده', 'کار', 'زمان', 'سال', 'ماه', 'روز',
        'اول', 'دوم', 'سوم', 'چهارم', 'پنجم', 'ششم', 'هفتم', 'هشتم', 'نهم', 'دهم',
        'فصل', 'بخش', 'مثال', 'تمرین', 'سوال', 'جواب', 'کتاب', 'درس', 'دانشگاه',
        'شرکت', 'مدیریت', 'پروژه', 'تیم', 'گروه', 'کاربر', 'سرور', 'شبکه',
    }
    
    @staticmethod
    def is_persian_char(char: str) -> bool:
        """Check if a character is Persian/Arabic."""
        if not char:
            return False
        code = ord(char)
        return (0x0600 <= code <= 0x06FF) or \
               (0x0750 <= code <= 0x077F) or \
               (0x08A0 <= code <= 0x08FF) or \
               (0xFB50 <= code <= 0xFDFF) or \
               (0xFE70 <= code <= 0xFEFF)
    
    @staticmethod
    def count_persian_words(text: str) -> int:
        """Count how many common Persian words are in text."""
        if not text:
            return 0
        words = set(text.split())
        return len(words & FileService.COMMON_PERSIAN_WORDS)
    
    @staticmethod
    def is_pdf_reversed(text: str) -> bool:
        """
        Detect if PDF text is reversed by comparing Persian word count
        in original vs reversed text.
        """
        if not text:
            return False
        
        # Count Persian words in original text
        original_count = FileService.count_persian_words(text)
        
        # Count Persian words in fully reversed text
        reversed_text = text[::-1]
        reversed_count = FileService.count_persian_words(reversed_text)
        
        logger.info(f"[PDF-DIRECTION] Original Persian words: {original_count}, Reversed: {reversed_count}")
        
        # If reversed version has more Persian words, PDF is reversed
        return reversed_count > original_count
    
    @staticmethod
    def fix_persian_pdf_text(text: str) -> str:
        """
        Fix Persian text extracted from PDF.
        Only reverse if text appears to be reversed.
        """
        if not text:
            return text
        
        # Check if text needs reversal
        if FileService.is_pdf_reversed(text):
            logger.info(f"[PDF-DIRECTION] Text appears REVERSED, fixing...")
            return text[::-1]
        else:
            logger.info(f"[PDF-DIRECTION] Text appears CORRECT, no fix needed")
            return text
    
    @staticmethod
    def clean_text(text: str) -> str:
        """Clean and normalize extracted text."""
        if not text:
            return ""
        
        # Remove null bytes and control characters
        text = re.sub(r'[\x00-\x08\x0b\x0c\x0e-\x1f\x7f]', '', text)
        
        # Fix Persian text direction if needed
        text = FileService.fix_persian_pdf_text(text)
        
        # Remove excessive whitespace
        text = re.sub(r'\n{3,}', '\n\n', text)
        text = re.sub(r' {2,}', ' ', text)
        
        # Strip each line
        lines = [line.strip() for line in text.split('\n')]
        
        # Remove empty lines at start and end
        while lines and not lines[0]:
            lines.pop(0)
        while lines and not lines[-1]:
            lines.pop()
        
        return '\n'.join(lines)
    
    @staticmethod
    def extract_text_from_pdf(file_path: str) -> str:
        """Extract text from PDF. Try pdfplumber first, then OCR if needed."""
        text = ""
        
        # Try pdfplumber first
        if HAS_PDFPLUMBER:
            try:
                with pdfplumber.open(file_path) as pdf:
                    for page in pdf.pages:
                        page_text = page.extract_text()
                        if page_text:
                            text += page_text + "\n"
                
                if text.strip():
                    logger.info(f"[PDF] Extracted {len(text)} chars with pdfplumber")
                    return text
                    
            except Exception as e:
                logger.warning(f"[PDF] pdfplumber failed: {e}")
        
        # Fallback to OCR for image-based PDFs
        if HAS_OCR:
            try:
                logger.info("[PDF] Trying OCR...")
                images = pdf2image.convert_from_path(file_path, dpi=300)
                
                for i, image in enumerate(images):
                    page_text = pytesseract.image_to_string(image, lang='eng+fas')
                    if page_text:
                        text += page_text + "\n"
                
                if text.strip():
                    logger.info(f"[PDF] Extracted {len(text)} chars with OCR")
                    return text
                    
            except Exception as e:
                logger.error(f"[PDF] OCR failed: {e}")
        
        return text
    
    @staticmethod
    def extract_text_from_docx(file_path: str) -> str:
        """Extract text from DOCX file."""
        if not HAS_DOCX:
            return "Error: python-docx not installed"
        
        try:
            doc = Document(file_path)
            text = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text += paragraph.text + "\n"
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                    if row_text:
                        text += row_text + "\n"
            
            logger.info(f"[DOCX] Extracted {len(text)} chars")
            return text
            
        except Exception as e:
            logger.error(f"[DOCX] Error: {e}")
            return f"Error reading DOCX: {str(e)}"
    
    @staticmethod
    def extract_text_from_txt(file_path: str) -> str:
        """Extract text from TXT file with encoding detection."""
        encodings = ['utf-8', 'utf-8-sig', 'latin-1', 'cp1256', 'iso-8859-1']
        
        for encoding in encodings:
            try:
                with open(file_path, 'r', encoding=encoding) as f:
                    text = f.read()
                logger.info(f"[TXT] Read with encoding: {encoding}")
                return text
            except (UnicodeDecodeError, UnicodeError):
                continue
        
        with open(file_path, 'r', encoding='utf-8', errors='ignore') as f:
            return f.read()
    
    @staticmethod
    def extract_text_from_pptx(file_path: str) -> str:
        """Extract text from PPTX file."""
        if not HAS_PPTX:
            return "Error: python-pptx not installed"
        
        try:
            prs = Presentation(file_path)
            text = ""
            
            for i, slide in enumerate(prs.slides, 1):
                slide_text = ""
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text += shape.text + "\n"
                
                if slide_text.strip():
                    text += f"--- Slide {i} ---\n{slide_text}\n"
            
            logger.info(f"[PPTX] Extracted {len(text)} chars")
            return text
            
        except Exception as e:
            logger.error(f"[PPTX] Error: {e}")
            return f"Error reading PPTX: {str(e)}"
    
    @staticmethod
    def extract_text_from_xlsx(file_path: str) -> str:
        """Extract text from XLSX file."""
        if not HAS_OPENPYXL:
            return "Error: openpyxl not installed"
        
        try:
            wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
            text = ""
            
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_text = f"--- Sheet: {sheet_name} ---\n"
                
                for row in ws.iter_rows(values_only=True):
                    row_text = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    if row_text.strip():
                        sheet_text += row_text + "\n"
                
                if sheet_text.strip():
                    text += sheet_text + "\n"
            
            wb.close()
            logger.info(f"[XLSX] Extracted {len(text)} chars")
            return text
            
        except Exception as e:
            logger.error(f"[XLSX] Error: {e}")
            return f"Error reading XLSX: {str(e)}"
    
    @staticmethod
    def extract_text_from_code(file_path: str) -> str:
        """Extract text from code files."""
        return FileService.extract_text_from_txt(file_path)
    
    @classmethod
    def extract_text(cls, file_path: str) -> str:
        """Extract text from any supported file format."""
        ext = os.path.splitext(file_path)[1].lower()
        
        extractors = {
            '.pdf': cls.extract_text_from_pdf,
            '.docx': cls.extract_text_from_docx,
            '.txt': cls.extract_text_from_txt,
            '.pptx': cls.extract_text_from_pptx,
            '.xlsx': cls.extract_text_from_xlsx,
            '.py': cls.extract_text_from_code,
            '.js': cls.extract_text_from_code,
            '.html': cls.extract_text_from_code,
            '.css': cls.extract_text_from_code,
            '.json': cls.extract_text_from_code,
            '.md': cls.extract_text_from_code,
        }
        
        extractor = extractors.get(ext)
        if extractor:
            text = extractor(file_path)
            return cls.clean_text(text)
        else:
            return f"Unsupported file format: {ext}"
